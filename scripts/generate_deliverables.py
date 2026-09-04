"""Generate the Review 1 presentation and a submission-ready workbook copy.

The presentation is intentionally generated from the repository's documented
implementation so it can be regenerated after each review.
"""
from pathlib import Path
from shutil import copy2

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "deliverables"
OUT.mkdir(exist_ok=True)

NAVY = RGBColor(10, 25, 55)
BLUE = RGBColor(27, 96, 170)
CYAN = RGBColor(55, 190, 205)
TEAL = RGBColor(24, 145, 145)
INK = RGBColor(22, 32, 51)
MUTED = RGBColor(89, 105, 128)
PALE = RGBColor(239, 247, 252)
WHITE = RGBColor(255, 255, 255)
ORANGE = RGBColor(239, 147, 61)
GREEN = RGBColor(41, 153, 104)


def text_box(slide, x, y, w, h, text, size=18, color=INK, bold=False,
             align=PP_ALIGN.LEFT, font="Aptos", margin=0.08):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def bullets(slide, x, y, w, h, items, size=17, color=INK, level_indent=0.22):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            value, level = item
        else:
            value, level = item, 0
        p.text = ("• " if level == 0 else "   – ") + value
        p.level = level
        p.font.name = "Aptos"
        p.font.size = Pt(size - level * 1)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return box


def rect(slide, x, y, w, h, fill, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line or fill
    return s


def header(slide, title, section):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = WHITE
    rect(slide, 0, 0, 13.333, 0.18, CYAN)
    text_box(slide, 0.65, 0.42, 9.9, 0.48, title, 27, NAVY, True)
    text_box(slide, 10.65, 0.48, 2.0, 0.28, section.upper(), 10, BLUE, True, PP_ALIGN.RIGHT)
    rect(slide, 0.65, 1.03, 12.0, 0.015, RGBColor(218, 229, 239))


def footer(slide, number):
    text_box(slide, 0.65, 7.15, 9.0, 0.2, "PolicyGPT  •  Project Work 2  •  Review 1", 9, MUTED)
    text_box(slide, 12.0, 7.15, 0.65, 0.2, str(number), 9, MUTED, align=PP_ALIGN.RIGHT)


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = NAVY
    rect(slide, 0, 0, 13.333, 0.22, CYAN)
    rect(slide, 0, 6.95, 13.333, 0.55, BLUE)
    text_box(slide, 0.75, 1.0, 11.8, 1.3, "PolicyGPT", 44, WHITE, True)
    text_box(slide, 0.78, 2.22, 11.4, 0.9, "An AI-Driven Policy Intelligence and Assistance Platform\nfor Government Schemes", 24, RGBColor(214, 238, 247), True)
    text_box(slide, 0.8, 3.48, 5.7, 0.52, "Project Work 2  |  Review 1", 20, CYAN, True)
    text_box(slide, 0.8, 4.25, 6.2, 1.0, "Guided by\nSanket S. Kulkarni\nAssistant Professor, Department of Machine Learning", 15, WHITE)
    text_box(slide, 7.4, 4.25, 5.0, 1.2, "Presented by\nSahana B.K.  •  1BM23AI162\nShamratha G.  •  1BM23AI173\nSuniksha Priya  •  1BM23AI192", 15, WHITE)
    text_box(slide, 0.8, 6.25, 6.0, 0.3, "B.M.S. College of Engineering  |  Department of Machine Learning", 11, RGBColor(214, 238, 247))
    text_box(slide, 9.0, 6.25, 3.4, 0.3, "Semester VII • Section C • 04 Sep 2026", 11, RGBColor(214, 238, 247), align=PP_ALIGN.RIGHT)


def add_slide(prs, title, section):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, title, section)
    footer(slide, len(prs.slides))
    return slide


def make_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    title_slide(prs)

    slide = add_slide(prs, "Agenda", "Review 1")
    agenda = ["01  Introduction and motivation", "02  Proposed system", "03  Implementation methodology", "04  System / module development", "05  Implementation progress", "06  Algorithm and model implementation", "07  Testing, validation and preliminary results", "08  Limitations, next steps and questions"]
    for i, item in enumerate(agenda):
        x = 0.9 if i < 4 else 6.8; y = 1.55 + (i % 4) * 1.15
        rect(slide, x, y, 5.55, 0.78, PALE, RGBColor(205, 224, 237), True)
        text_box(slide, x + 0.18, y + 0.18, 5.15, 0.35, item, 18, NAVY, True)

    slide = add_slide(prs, "Why PolicyGPT?", "01  Introduction")
    rect(slide, 0.8, 1.4, 3.65, 4.95, NAVY, NAVY, True)
    text_box(slide, 1.1, 1.8, 3.05, 0.5, "The access gap", 23, WHITE, True)
    text_box(slide, 1.1, 2.48, 3.0, 2.85, "Government schemes are available, but often\n\n• scattered across portals\n• buried in long PDFs\n• written in legal terminology\n• hard to compare or personalize", 17, RGBColor(220, 237, 247))
    text_box(slide, 4.95, 1.5, 7.55, 0.6, "Two users, one evidence base", 23, NAVY, True)
    rect(slide, 4.95, 2.25, 3.45, 3.5, RGBColor(237, 249, 246), RGBColor(177, 224, 210), True)
    text_box(slide, 5.25, 2.6, 2.8, 0.45, "Citizen assistant", 20, TEAL, True)
    bullets(slide, 5.2, 3.18, 2.95, 2.1, ["Plain-language explanations", "Likely eligibility checks", "Application steps and sources"], 16)
    rect(slide, 8.75, 2.25, 3.45, 3.5, RGBColor(246, 242, 252), RGBColor(211, 196, 234), True)
    text_box(slide, 9.05, 2.6, 2.8, 0.45, "Policy intelligence", 20, BLUE, True)
    bullets(slide, 9.0, 3.18, 2.95, 2.1, ["Compare policies", "Find gaps and overlaps", "Export auditable reports"], 16)
    text_box(slide, 4.95, 6.05, 7.25, 0.35, "Responsible-AI principle: show evidence, uncertainty and issuing-agency responsibility.", 14, ORANGE, True)

    slide = add_slide(prs, "Proposed system and architecture", "02  System")
    text_box(slide, 0.8, 1.28, 11.6, 0.45, "One traceable pipeline serves both the citizen and analyst workflows.", 17, MUTED)
    nodes = [("Browser UI", 0.85, BLUE), ("FastAPI\nAPI", 2.55, TEAL), ("Orchestrator", 4.25, NAVY), ("Domain\nrouter", 6.0, BLUE), ("Hybrid\nretrieval", 7.72, TEAL), ("Generate +\nvalidate", 9.45, NAVY), ("Report +\nhistory", 11.12, BLUE)]
    for label, x, fill in nodes:
        rect(slide, x, 2.35, 1.35, 1.0, fill, fill, True)
        text_box(slide, x + 0.08, 2.62, 1.19, 0.45, label, 14, WHITE, True, PP_ALIGN.CENTER)
    for i in range(len(nodes)-1):
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(nodes[i][1]+1.35), Inches(2.85), Inches(nodes[i+1][1]), Inches(2.85))
        line.line.color.rgb = CYAN; line.line.width = Pt(2)
    rect(slide, 1.05, 4.25, 5.3, 1.35, PALE, RGBColor(205, 224, 237), True)
    text_box(slide, 1.3, 4.48, 4.8, 0.8, "Citizen path\nquestion → explanation → eligibility → application guidance", 16, NAVY, True)
    rect(slide, 6.85, 4.25, 5.3, 1.35, RGBColor(246, 242, 252), RGBColor(211, 196, 234), True)
    text_box(slide, 7.1, 4.48, 4.8, 0.8, "Analyst path\nretrieve → compare → identify gaps/overlaps → export report", 16, NAVY, True)
    text_box(slide, 1.05, 6.15, 11.1, 0.3, "SQLite stores sessions, messages, metadata and generated reports; source metadata travels with every chunk.", 14, MUTED)

    slide = add_slide(prs, "Implementation methodology", "03  Methodology")
    steps = [("1", "Ingest", "clean policy text + source metadata"), ("2", "Route", "domain and intent classification"), ("3", "Retrieve", "lexical + semantic ranking"), ("4", "Generate", "role-aware grounded answer"), ("5", "Validate", "claim audit + confidence"), ("6", "Report", "citations, warnings, export")]
    for i, (num, title, desc) in enumerate(steps):
        x = 0.8 + (i % 3) * 4.1; y = 1.45 + (i // 3) * 2.15
        rect(slide, x, y, 3.55, 1.45, PALE if i < 3 else RGBColor(241, 247, 244), RGBColor(205, 224, 237), True)
        rect(slide, x + 0.2, y + 0.22, 0.52, 0.52, BLUE if i < 3 else TEAL, None, True)
        text_box(slide, x + 0.2, y + 0.3, 0.52, 0.28, num, 16, WHITE, True, PP_ALIGN.CENTER)
        text_box(slide, x + 0.9, y + 0.2, 2.35, 0.35, title, 18, NAVY, True)
        text_box(slide, x + 0.9, y + 0.7, 2.35, 0.5, desc, 13, MUTED)
    text_box(slide, 0.85, 6.2, 11.4, 0.45, "Hybrid score = 0.35 lexical relevance + 0.65 deterministic semantic bigram similarity. A bounded repair pass runs when confidence < 0.75 or hallucination risk > 0.25.", 14, ORANGE, True)

    slide = add_slide(prs, "System and module development", "04  Modules")
    modules = [("Policy corpus", "Six domains: education, agriculture, MSME, finance, health, disaster", GREEN), ("Citizen assistant", "Question answering, simplified explanations, profile-aware checks", BLUE), ("Policy intelligence", "Retrieval, comparison, gap/overlap analysis and report export", TEAL), ("Trust layer", "Citations, claim audits, confidence and hallucination warnings", ORANGE), ("Platform", "FastAPI, SQLite, browser UI, Docker and session history", NAVY)]
    for i,(name,desc,fill) in enumerate(modules):
        y=1.35+i*1.02
        rect(slide, 0.9, y, 2.45, 0.72, fill, fill, True); text_box(slide, 1.05, y+0.2, 2.15, 0.3, name, 16, WHITE, True)
        rect(slide, 3.65, y, 8.55, 0.72, PALE, RGBColor(216, 229, 238), True); text_box(slide, 3.9, y+0.18, 8.0, 0.35, desc, 15, INK)
    text_box(slide, 0.95, 6.55, 11.2, 0.3, "Implementation status: prototype path is functional end-to-end; seed records remain illustrative until replaced with approved archived sources.", 14, ORANGE, True)

    slide = add_slide(prs, "Implementation progress", "05  Progress")
    weeks = [("W1–2", "Architecture, backlog, FastAPI, SQLite, UI and Docker"), ("W3–4", "Routing, hybrid retrieval, generation, eligibility and validation"), ("W5–6", "Bounded repair, Markdown reports, source links and history"), ("W7–8", "API/UI smoke tests, presentation evidence and next-phase plan")]
    for i,(wk,desc) in enumerate(weeks):
        x=0.9+i*3.05
        rect(slide, x, 1.65, 2.55, 3.45, NAVY if i==3 else PALE, NAVY if i==3 else RGBColor(205,224,237), True)
        text_box(slide, x+0.22, 1.98, 2.1, 0.45, wk, 24, CYAN if i==3 else BLUE, True)
        text_box(slide, x+0.22, 2.75, 2.08, 1.75, desc, 16, WHITE if i==3 else INK, True if i==3 else False)
        text_box(slide, x+0.22, 4.55, 2.1, 0.28, "completed", 11, RGBColor(174, 224, 199) if i==3 else GREEN, True)
    text_box(slide, 0.95, 5.75, 11.2, 0.7, "Current review baseline\nA runnable prototype, documented technical decisions, deterministic fallback, test suite, and demo-ready workflow.", 17, NAVY, True)

    slide = add_slide(prs, "Algorithm and model implementation", "06  Retrieval")
    rect(slide, 0.85, 1.45, 5.65, 4.65, PALE, RGBColor(205,224,237), True)
    text_box(slide, 1.15, 1.78, 4.9, 0.4, "Hybrid retrieval", 23, NAVY, True)
    bullets(slide, 1.1, 2.45, 5.0, 2.8, ["Lexical signal preserves exact scheme terms and legal phrases.", "Deterministic semantic bigrams capture paraphrases without an external model dependency.", "Weighted ranking: 35% lexical + 65% semantic.", "Top chunks retain source, page and domain metadata."], 16)
    rect(slide, 6.8, 1.45, 5.65, 4.65, RGBColor(246,242,252), RGBColor(211,196,234), True)
    text_box(slide, 7.1, 1.78, 4.9, 0.4, "Validation loop", 23, NAVY, True)
    bullets(slide, 7.05, 2.45, 5.0, 2.8, ["Generation is constrained by retrieved evidence.", "Claims are matched back to evidence and citations.", "Confidence and hallucination risk are surfaced to the user.", "One bounded repair pass broadens context for weak answers."], 16)
    text_box(slide, 0.95, 6.35, 11.4, 0.35, "Design choice: deterministic baseline makes the prototype reproducible and offline-capable; production can add approved embedding models.", 14, ORANGE, True)

    slide = add_slide(prs, "Testing, validation and preliminary results", "07  Evidence")
    checks = [("Retrieval", "ranking + repeatability tests", GREEN), ("API", "health, query and export smoke tests", BLUE), ("Trust", "citation presence and claim audit checks", TEAL), ("Review", "workbook, slide and limitation checklist", ORANGE)]
    for i,(name,desc,fill) in enumerate(checks):
        y=1.42+i*1.0
        rect(slide, 0.95, y, 2.25, 0.68, fill, fill, True); text_box(slide, 1.15, y+0.19, 1.85, 0.3, name, 17, WHITE, True)
        rect(slide, 3.55, y, 8.55, 0.68, PALE, RGBColor(216,229,238), True); text_box(slide, 3.8, y+0.18, 8.0, 0.3, desc, 15, INK)
    rect(slide, 0.95, 5.65, 11.15, 0.82, RGBColor(237,249,246), RGBColor(177,224,210), True)
    text_box(slide, 1.2, 5.87, 10.65, 0.38, "Preliminary result: one request returns a grounded answer, source links, validation metrics and an exportable report. Latency is corpus/provider dependent; offline fallback has no external API dependency.", 15, TEAL, True)

    slide = add_slide(prs, "Limitations, next steps and team contribution", "08  Close")
    rect(slide, 0.85, 1.35, 5.6, 4.9, RGBColor(255,247,238), RGBColor(244,214,177), True)
    text_box(slide, 1.15, 1.68, 4.8, 0.4, "Next phase", 22, ORANGE, True)
    bullets(slide, 1.1, 2.3, 4.9, 3.3, ["Replace illustrative records with approved, dated official documents.", "Add multilingual support, authentication and richer analyst views.", "Benchmark retrieval, citation correctness, eligibility accuracy and latency.", "Add human-in-the-loop review and user evaluation."], 16)
    rect(slide, 6.8, 1.35, 5.6, 4.9, PALE, RGBColor(205,224,237), True)
    text_box(slide, 7.1, 1.68, 4.8, 0.4, "Team contribution", 22, BLUE, True)
    bullets(slide, 7.05, 2.3, 4.9, 3.3, ["Sahana — corpus and API contracts", "Shamratha — retrieval and orchestration", "Suniksha — UI, testing and documentation", "All members — live query demo and technical Q&A"], 16)
    text_box(slide, 0.95, 6.5, 11.2, 0.3, "Questions: Why hybrid retrieval? Why validate claims? Why not replace official portals?", 15, NAVY, True, PP_ALIGN.CENTER)

    slide = add_slide(prs, "Suggestions / Questions", "Discussion")
    text_box(slide, 1.05, 2.1, 11.2, 1.0, "Thank you.\nWe welcome feedback on corpus scope, evaluation design and deployment priorities.", 28, NAVY, True, PP_ALIGN.CENTER)
    text_box(slide, 1.05, 4.3, 11.2, 0.5, "PolicyGPT is a support tool; official departments remain the source of truth for eligibility and application decisions.", 15, ORANGE, True, PP_ALIGN.CENTER)

    slide = add_slide(prs, "Thank you", "Close")
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = NAVY
    rect(slide, 0, 0, 13.333, 0.22, CYAN)
    text_box(slide, 0.8, 2.15, 11.8, 0.75, "Thank you", 42, WHITE, True, PP_ALIGN.CENTER)
    text_box(slide, 0.8, 3.25, 11.8, 0.65, "PolicyGPT  •  Review 1  •  Project Work 2", 22, RGBColor(214,238,247), True, PP_ALIGN.CENTER)
    text_box(slide, 0.8, 5.4, 11.8, 0.4, "Sahana B.K.  |  Shamratha G.  |  Suniksha Priya", 16, CYAN, True, PP_ALIGN.CENTER)

    target = OUT / "PolicyGPT_Project_Work_2_Review_1.pptx"
    prs.save(target)
    return target


if __name__ == "__main__":
    pptx = make_deck()
    workbook = ROOT / "docs" / "weekly-progress-report-filled.xlsx"
    copy2(workbook, OUT / "PolicyGPT_Weekly_Progress_Report_Review_1.xlsx")
    print(pptx)
    print(OUT / "PolicyGPT_Weekly_Progress_Report_Review_1.xlsx")
